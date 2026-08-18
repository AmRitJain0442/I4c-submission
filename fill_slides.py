"""Fill the official KLA idea-submission template with our content + results.

Usage: python fill_slides.py [--psno PS01]
Reads:  docs/template.pptx, docs/results.json (optional), docs/figures/*.png
Writes: docs/gitcommitall_<psno>.pptx
"""
import argparse
import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn

REPO_URL = "https://github.com/AmRitJain0442/I4c-submission"
TEAM = "gitcommitall"
MEMBERS = ["Amrit Lahari", "Shubham Gangwar", "Saiswaroop Natarajan"]
COLLEGE = "BITS Pilani"
EMAIL = "f20230442@pilani.bits-pilani.ac.in"

INK = RGBColor(0x0B, 0x0B, 0x0B)


def set_text(sh, text, anchor_top=False):
    tf = sh.text_frame
    if anchor_top:
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
    p = tf.paragraphs[0]
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    r.text = text
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape {shape_id} not on slide")


def clone_slide(prs, src):
    """Clone a slide's shapes into a new slide, remapping relationship ids."""
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):
        shp._element.getparent().remove(shp._element)
    rel_attrs = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in rel_attrs:
                rid = node.get(attr)
                if not rid:
                    continue
                try:
                    rel = src.part.rels[rid]
                except KeyError:
                    continue
                if rel.is_external:
                    new_rid = dst.part.rels._add_relationship(rel.reltype, rel.target_ref, True)
                else:
                    new_rid = dst.part.rels._add_relationship(rel.reltype, rel.target_part, False)
                node.set(attr, new_rid)
        dst.shapes._spTree.append(el)
    return dst


def move_slide(prs, old_index, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[old_index]
    lst.remove(el)
    lst.insert(new_index, el)


def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    el = list(lst)[index]
    rId = el.get(qn("r:id"))
    prs.part.drop_rel(rId)
    lst.remove(el)


def fmt(v, spec):
    return spec.format(v) if v is not None else "TBD"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--psno", default="PS")
    args = ap.parse_args()

    res = {}
    if Path("docs/results.json").exists():
        res = json.loads(Path("docs/results.json").read_text())
    baseline = res.get("baseline", {})
    models = res.get("models", {})
    naf, swin = models.get("nafnet", {}), models.get("swinir", {})
    winner_name = res.get("winner", "TBD")
    win = models.get(winner_name, {})

    prs = Presentation("docs/template.pptx")
    s = prs.slides

    # --- slide 1: team details ---
    t = s[1]
    set_text(by_id(t, 54), TEAM)
    for sid, name in zip((55, 57, 59), MEMBERS):
        set_text(by_id(t, sid), name)
    set_text(by_id(t, 61), "-")
    for sid in (56, 58, 60, 62):
        set_text(by_id(t, sid), "")
    set_text(by_id(t, 63), COLLEGE)
    set_text(by_id(t, 65), EMAIL)

    # --- slide 2: problem statement ---
    p = s[2]
    set_text(by_id(p, 16), "AI-Based Restoration of Degraded Images for Semiconductor Inspection (KLA)")
    set_text(by_id(p, 19),
             "Microscopic inspection images in semiconductor fabs are degraded by three simultaneous "
             "modes: speckle noise (pixel values overshoot the valid [0,1] range - we measured inputs "
             "spanning -0.21 to 1.87), additive Gaussian noise, and 2x loss of spatial resolution "
             "(256x256 -> 128x128). Degradation hides the fine structures and defects that inspection "
             "depends on. The task: a single AI model that jointly denoises and super-resolves each "
             "128x128 input back to a clean 256x256 image, scored on SSIM, PSNR, LPIPS, inference "
             "speed (H100) and generalization to out-of-distribution structures.", anchor_top=True)

    # --- slide 3: idea description ---
    i = s[3]
    set_text(by_id(i, 16),
             "One network, one pass: joint denoising + 2x super-resolution, trained end-to-end on the "
             "3,200 paired images provided by KLA. We benchmarked a CNN track and a transformer track "
             "under an identical training harness and ship the empirically better model.")
    set_text(by_id(i, 19),
             "Learn only the correction: the network output is added to a bicubic 2x upsample of the "
             "input (global residual), so capacity is spent on removing noise and restoring detail rather "
             "than rebuilding the image. Charbonnier loss, cosine LR, AMP, flip/rotation augmentation.", anchor_top=True)
    set_text(by_id(i, 22),
             "Restores fine defect-relevant structure without hallucination: a compact restoration "
             "network (NAFNet-SR 1.13M params / SwinIR 5.16M params) maps each degraded 128x128 frame "
             "to a clean 256x256 frame in a single forward pass, running in milliseconds per image.", anchor_top=True)

    # --- slide 4: proposed solution ---
    ps_ = s[4]
    set_text(by_id(ps_, 16),
             "Two architectures trained and compared on a clean held-out validation split of 400 images; "
             "the better one (by PSNR / SSIM / LPIPS and latency) is submitted.")
    set_text(by_id(ps_, 19),
             "Track A - NAFNet-SR: nonlinear-activation-free blocks (LayerNorm, SimpleGate, simplified "
             "channel attention) at low resolution, PixelShuffle 2x head, bicubic global residual. "
             "Track B - SwinIR: windowed self-attention (window 8) that suits the repeating patterns of "
             "semiconductor structures. Both trained on random 64x64 crops, batch 32, ~2h on GCP "
             "(A100 40GB + L4). Validation every 5 epochs; best checkpoint kept. Optional LPIPS "
             "fine-tune and x8 self-ensemble (TTA) evaluated against the latency budget.", anchor_top=True)

    # --- slide 5: innovation ---
    inn = s[5]
    set_text(by_id(inn, 16),
             "Engineering rigor over exotic architecture: measured model selection, leakage-free "
             "validation, and speed-aware design choices.")
    set_text(by_id(inn, 19),
             "Global bicubic-residual design: the network learns only the degradation correction, which "
             "stabilizes training and preserves true structure (no hallucinated geometry - critical for "
             "defect inspection). Range-aware handling keeps speckle overshoot information ([-0.2, 1.9]) "
             "instead of clipping the input.", anchor_top=True)
    set_text(by_id(inn, 22),
             "CNN-vs-transformer bake-off under one harness, selected on measured PSNR / SSIM / LPIPS "
             "and per-image latency rather than paper claims. A strict 400-image held-out validation "
             "split (never trained on) keeps model selection honest, and the fully-provided degradation "
             "pairs mean no synthetic-degradation mismatch.", anchor_top=True)

    # --- slide 6: impact ---
    imp = s[6]
    set_text(by_id(imp, 16),
             "Sharper inspection images directly improve defect detection reliability, reduce false "
             "rejects, and let fabs run faster or cheaper imaging while software recovers the quality.")
    set_text(by_id(imp, 20),
             "Restores inspection-grade image quality from noisy, low-resolution captures in a single "
             "millisecond-scale forward pass per frame - fits inline inspection throughput.")
    set_text(by_id(imp, 24),
             f"+{fmt(win.get('psnr') - baseline.get('psnr') if win.get('psnr') and baseline.get('psnr') else None, '{:.1f}')} dB PSNR over bicubic "
             f"({fmt(baseline.get('psnr'), '{:.2f}')} -> {fmt(win.get('psnr'), '{:.2f}')} dB); "
             f"SSIM {fmt(baseline.get('ssim'), '{:.3f}')} -> {fmt(win.get('ssim'), '{:.3f}')}; "
             f"LPIPS {fmt(win.get('lpips'), '{:.3f}')}; {fmt(win.get('ms'), '{:.0f}')} ms/image.")

    # --- results slide (clone of slide 2 chrome) ---
    r = clone_slide(prs, s[2])
    set_text(by_id(r, 8), "Results - Quantitative & Visual")
    for sid in (14, 15, 16, 17, 18, 19):
        try:
            sh = by_id(r, sid)
            sh._element.getparent().remove(sh._element)
        except KeyError:
            pass

    rows = [
        ("Model", "PSNR (dB)", "SSIM", "LPIPS", "ms / image"),
        ("Bicubic x2 (baseline)", fmt(baseline.get("psnr"), "{:.2f}"), fmt(baseline.get("ssim"), "{:.3f}"),
         fmt(baseline.get("lpips"), "{:.3f}"), "-"),
        ("NAFNet-SR (1.13M)", fmt(naf.get("psnr"), "{:.2f}"), fmt(naf.get("ssim"), "{:.3f}"),
         fmt(naf.get("lpips"), "{:.3f}"), fmt(naf.get("ms"), "{:.1f}")),
        ("SwinIR (5.16M)", fmt(swin.get("psnr"), "{:.2f}"), fmt(swin.get("ssim"), "{:.3f}"),
         fmt(swin.get("lpips"), "{:.3f}"), fmt(swin.get("ms"), "{:.1f}")),
    ]
    tbl = r.shapes.add_table(len(rows), 5, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6)).table
    tbl.columns[0].width = Inches(3.9)
    for ci in range(1, 5):
        tbl.columns[ci].width = Inches(1.9)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(12 if ri == 0 else 11)
            para.runs[0].font.bold = ri == 0 or (res.get("winner") == "nafnet" and ri == 2) \
                or (res.get("winner") == "swinir" and ri == 3)

    figs = Path("docs/figures")
    if (figs / "training_curves.png").exists():
        r.shapes.add_picture(str(figs / "training_curves.png"), Inches(0.9), Inches(4.15), width=Inches(5.6))
    comp = sorted(figs.glob("compare_*.png"))
    if comp:
        r.shapes.add_picture(str(comp[0]), Inches(6.7), Inches(4.3), width=Inches(5.7))
    if (figs / "metrics_panels.png").exists() and not comp:
        r.shapes.add_picture(str(figs / "metrics_panels.png"), Inches(6.7), Inches(4.3), width=Inches(5.7))

    move_slide(prs, len(prs.slides._sldIdLst) - 1, 7)

    # --- slide idx shifts: original 7,8,9 are now 8,9,10 ---
    tech = s[8]
    set_text(by_id(tech, 16),
             "PyTorch 2.x (AMP, torch.compile), timm, scikit-image, LPIPS; trained on GCP (A100 40GB + "
             "L4) in under 2 hours per model - cheap to retrain as processes drift.")
    set_text(by_id(tech, 19),
             "Deliverable-first design: a standalone evaluate.py restores a directory of .npy images "
             "and reports PSNR / SSIM / LPIPS + mean per-image latency with zero manual editing "
             "(KLA benchmark-ready). Resumable training, pinned requirements.txt, and weights shipped "
             "in-repo make the full pipeline reproducible end to end.", anchor_top=True)

    gh = s[9]
    set_text(by_id(gh, 20), REPO_URL)

    ref = s[10]
    set_text(by_id(ref, 19),
             "Residual and attention-based image restoration; joint denoising + super-resolution on "
             "paired degraded/clean data; perceptual metrics (LPIPS) alongside PSNR/SSIM for "
             "inspection-relevant fidelity.", anchor_top=True)
    for sid, txt in ((26, "Chen et al., Simple Baselines for Image Restoration (NAFNet), ECCV 2022 - arXiv:2204.04676"),
                     (28, "Liang et al., SwinIR: Image Restoration Using Swin Transformer, ICCVW 2021 - arXiv:2108.10257"),
                     (29, "Zhang et al., The Unreasonable Effectiveness of Deep Features (LPIPS), CVPR 2018 - arXiv:1801.03924")):
        sh = by_id(ref, sid)
        sh.width = Inches(10.0)
        set_text(sh, txt)

    # --- drop instructions slide ---
    delete_slide(prs, 0)

    out = f"docs/{TEAM}_{args.psno}.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
